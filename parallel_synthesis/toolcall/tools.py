import csv
import io
import json
import mimetypes
import os
import re
import shutil
import subprocess
import tempfile
import time
import urllib.parse
import urllib.request
import xml.etree.ElementTree as ET
import zipfile
from html import unescape
from pathlib import Path
from typing import Any, Dict, List, Optional


def _http_get(url: str, timeout: int = 25, headers: Optional[Dict[str, str]] = None) -> str:
    req_headers = {
        "User-Agent": os.getenv(
            "TOOL_HTTP_USER_AGENT",
            "Mozilla/5.0 (X11; Linux x86_64) AppleWebKit/537.36 "
            "(KHTML, like Gecko) Chrome/124.0.0.0 Safari/537.36",
        ),
        "Accept-Language": "en-US,en;q=0.9",
    }
    if headers:
        req_headers.update(headers)
    req = urllib.request.Request(url, headers=req_headers)
    with urllib.request.urlopen(req, timeout=timeout) as resp:
        raw = resp.read()
    return raw.decode("utf-8", errors="ignore")


def _strip_html(text: str) -> str:
    text = re.sub(r"<script[\\s\\S]*?</script>", " ", text, flags=re.IGNORECASE)
    text = re.sub(r"<style[\\s\\S]*?</style>", " ", text, flags=re.IGNORECASE)
    text = re.sub(r"<[^>]+>", " ", text)
    text = unescape(text)
    text = re.sub(r"\\s+", " ", text)
    return text.strip()


def _is_duckduckgo_bot_challenge(html: str) -> bool:
    lowered = html.lower()
    markers = (
        "unfortunately, bots use duckduckgo too",
        "please complete the following challenge",
        "select all squares containing",
        "duckduckgo.com/anomaly.js",
        "ddg-human-verification",
    )
    return any(marker in lowered for marker in markers)


class ToolSpec:
    def __init__(self, name: str, description: str, parameters: Dict[str, Any]) -> None:
        self.name = name
        self.description = description
        self.parameters = parameters


class BaseTool:
    spec: ToolSpec

    def call(self, arguments: Dict[str, Any]) -> str:
        raise NotImplementedError


class SearchTool(BaseTool):
    def __init__(self) -> None:
        self.spec = ToolSpec(
            name="search",
            description=(
                "Perform web search. Accepts one query string or a list of queries; "
                "returns top snippets and URLs for each query."
            ),
            parameters={
                "type": "object",
                "properties": {
                    "query": {
                        "type": ["string", "array"],
                        "items": {"type": "string"},
                        "description": "Search query or list of queries.",
                    }
                },
                "required": ["query"],
            },
        )

    def _search_serper(self, query: str) -> str:
        api_key = os.getenv("SERPER_KEY_ID", "").strip()
        if not api_key:
            raise RuntimeError("SERPER_KEY_ID is not set")
        payload = json.dumps({"q": query}).encode("utf-8")
        req = urllib.request.Request(
            "https://google.serper.dev/search",
            data=payload,
            headers={
                "X-API-KEY": api_key,
                "Content-Type": "application/json",
            },
            method="POST",
        )
        with urllib.request.urlopen(req, timeout=25) as resp:
            row = json.loads(resp.read().decode("utf-8", errors="ignore"))

        results = row.get("organic", [])
        chunks = []
        for i, r in enumerate(results[:8], start=1):
            title = str(r.get("title", "")).strip()
            link = str(r.get("link", "")).strip()
            snippet = str(r.get("snippet", "")).strip()
            chunks.append(f"{i}. {title}\nURL: {link}\nSnippet: {snippet}")
        if not chunks:
            return f"No search results found for: {query}"
        return f"Search results for '{query}':\n" + "\n\n".join(chunks)

    def _search_duckduckgo(self, query: str) -> str:
        url = "https://duckduckgo.com/html/?" + urllib.parse.urlencode({"q": query})
        html = _http_get(url, timeout=25)
        if _is_duckduckgo_bot_challenge(html):
            return (
                "[search] DuckDuckGo blocked the automated fallback with an anti-bot "
                f"challenge for query: {query}. Configure SERPER_KEY_ID for API-backed "
                "search; the HTML fallback is unreliable for benchmark runs."
            )

        entries = []
        pattern = re.compile(
            r'<a[^>]*class="result__a"[^>]*href="([^"]+)"[^>]*>(.*?)</a>[\\s\\S]{0,600}?'
            r'<a[^>]*class="result__snippet"[^>]*>(.*?)</a>',
            flags=re.IGNORECASE,
        )
        for i, m in enumerate(pattern.finditer(html), start=1):
            href = unescape(m.group(1))
            title = _strip_html(m.group(2))
            snippet = _strip_html(m.group(3))
            entries.append(f"{i}. {title}\nURL: {href}\nSnippet: {snippet}")
            if i >= 8:
                break

        if not entries:
            text = _strip_html(html)
            return (
                f"[search] Search fallback returned no parseable results for query: {query}. "
                f"Preview: {text[:500]}"
            )
        return f"Search results for '{query}':\n" + "\n\n".join(entries)

    def _search_one(self, query: str) -> str:
        try:
            return self._search_serper(query)
        except Exception:
            return self._search_duckduckgo(query)

    def call(self, arguments: Dict[str, Any]) -> str:
        raw = arguments.get("query")
        if isinstance(raw, str):
            queries = [raw]
        elif isinstance(raw, list):
            queries = [str(x) for x in raw if str(x).strip()]
        else:
            return "[search] invalid arguments: expected query as string or list"

        outputs = [self._search_one(q.strip()) for q in queries[:6]]
        return "\n\n=======\n\n".join(outputs)


class VisitTool(BaseTool):
    def __init__(self) -> None:
        self._jina_keys = [
            x.strip() for x in os.getenv("JINA_API_KEYS", "").split(",") if x.strip()
        ]
        self._jina_timeout = int(os.getenv("JINA_VISIT_TIMEOUT", "35"))
        self._jina_retries = int(os.getenv("JINA_VISIT_RETRIES", "3"))
        self._direct_timeout = int(os.getenv("DIRECT_VISIT_TIMEOUT", "20"))
        self._max_chars = int(os.getenv("VISIT_MAX_CHARS", "18000"))

        self.spec = ToolSpec(
            name="visit",
            description="Visit webpage(s) and return extracted text relevant to a goal.",
            parameters={
                "type": "object",
                "properties": {
                    "url": {
                        "type": ["string", "array"],
                        "items": {"type": "string"},
                        "description": "URL or list of URLs",
                    },
                    "goal": {
                        "type": "string",
                        "description": "What information to extract from the page",
                    },
                },
                "required": ["url", "goal"],
            },
        )

    @staticmethod
    def _normalize_web_url(url: str) -> str:
        s = str(url).strip()
        if not s:
            return s
        if s.startswith(("http://", "https://")):
            return s
        return "https://" + s

    def _jina_fetch_once(self, target_url: str, api_key: str = "") -> str:
        jina_url = f"https://r.jina.ai/{target_url}"
        headers: Dict[str, str] = {}
        if api_key:
            headers["Authorization"] = f"Bearer {api_key}"
        return _http_get(jina_url, timeout=self._jina_timeout, headers=headers)

    def _read_via_jina(self, target_url: str) -> str:
        """
        Try authenticated Jina first (if keys exist), then anonymous Jina.
        Retries and rotates API keys to improve robustness.
        """
        keys = self._jina_keys[:] if self._jina_keys else [""]
        last_error = ""
        for attempt in range(max(1, self._jina_retries)):
            for key in keys:
                try:
                    txt = self._jina_fetch_once(target_url, api_key=key).strip()
                    if txt:
                        return txt
                except Exception as exc:
                    last_error = str(exc)
            # small backoff to reduce transient throttling issues
            time.sleep(min(0.5 * (attempt + 1), 1.5))
        if last_error:
            raise RuntimeError(last_error)
        raise RuntimeError("empty response from Jina reader")

    def _read_via_direct_http(self, target_url: str) -> str:
        # First try direct URL, then "www." variant as fallback.
        candidates = [target_url]
        parsed = urllib.parse.urlparse(target_url)
        if parsed.scheme and parsed.netloc and not parsed.netloc.startswith("www."):
            with_www = parsed._replace(netloc="www." + parsed.netloc)
            candidates.append(urllib.parse.urlunparse(with_www))

        last_error = ""
        for cand in candidates:
            try:
                html = _http_get(cand, timeout=self._direct_timeout)
                return _strip_html(html)
            except Exception as exc:
                last_error = str(exc)
        raise RuntimeError(last_error or "direct fetch failed")

    def _read_url(self, url: str) -> str:
        url = str(url).strip()
        if not url:
            return "[visit] empty URL"

        if os.path.exists(url):
            try:
                with open(url, "r", encoding="utf-8", errors="ignore") as fh:
                    txt = fh.read()
                return txt[:12000]
            except Exception as exc:
                return f"[visit] failed to read local file {url}: {exc}"

        target_url = self._normalize_web_url(url)
        try:
            txt = self._read_via_jina(target_url)
            return txt[: self._max_chars]
        except Exception as jina_exc:
            jina_err = str(jina_exc)

        try:
            txt = self._read_via_direct_http(target_url)
            return txt[: self._max_chars]
        except Exception as direct_exc:
            direct_err = str(direct_exc)
            # Provide compact diagnostics to help the model decide next action.
            return (
                f"[visit] failed to fetch URL {url}\n"
                f"- jina_error: {jina_err}\n"
                f"- direct_error: {direct_err}"
            )

    def call(self, arguments: Dict[str, Any]) -> str:
        raw_urls = arguments.get("url")
        goal = str(arguments.get("goal", "")).strip()
        if isinstance(raw_urls, str):
            urls = [raw_urls]
        elif isinstance(raw_urls, list):
            urls = [str(u) for u in raw_urls]
        else:
            return "[visit] invalid arguments: url must be string or list"

        outputs = []
        for u in urls[:5]:
            content = self._read_url(u)
            outputs.append(
                f"URL: {u}\nGoal: {goal}\nExtracted content:\n{content}"
            )
        return "\n\n=======\n\n".join(outputs)


class PythonInterpreterTool(BaseTool):
    def __init__(self) -> None:
        self.spec = ToolSpec(
            name="PythonInterpreter",
            description=(
                "Execute Python code. Use <code>...</code> inside a tool call. "
                "Any output you want must be printed."
            ),
            parameters={
                "type": "object",
                "properties": {
                    "code": {"type": "string", "description": "Python source code"}
                },
                "required": ["code"],
            },
        )

    def call(self, arguments: Dict[str, Any]) -> str:
        code = str(arguments.get("code", "")).strip()
        if not code:
            return "[PythonInterpreter] empty code"

        timeout_sec = int(os.getenv("TOOL_PY_TIMEOUT", "35"))
        with tempfile.NamedTemporaryFile("w", suffix=".py", delete=False, encoding="utf-8") as tmp:
            tmp.write(code)
            tmp_path = tmp.name

        try:
            proc = subprocess.run(
                ["python", tmp_path],
                capture_output=True,
                text=True,
                timeout=timeout_sec,
            )
            out = proc.stdout.strip()
            err = proc.stderr.strip()
            parts = []
            if out:
                parts.append("stdout:\n" + out)
            if err:
                parts.append("stderr:\n" + err)
            if not parts:
                parts.append("Finished execution.")
            return "\n\n".join(parts)
        except subprocess.TimeoutExpired:
            return f"[PythonInterpreter] timeout after {timeout_sec}s"
        except Exception as exc:
            return f"[PythonInterpreter] execution failed: {exc}"
        finally:
            try:
                os.remove(tmp_path)
            except OSError:
                pass


class ParseFileTool(BaseTool):
    def __init__(self, file_root: str = "") -> None:
        self.file_root = file_root.strip()
        self.max_chars = int(os.getenv("PARSE_FILE_MAX_CHARS", "30000"))
        self.spec = ToolSpec(
            name="parse_file",
            description=(
                "Parse local files and return readable text. Supports common formats "
                "including txt/md/json/csv/tsv/html/xml/pdf/docx/pptx/xlsx/xls/zip/images. "
                "Good for attached GAIA files."
            ),
            parameters={
                "type": "object",
                "properties": {
                    "files": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "File paths",
                    }
                },
                "required": ["files"],
            },
        )

    def _resolve(self, path: str) -> str:
        p = str(path).strip()
        if not p:
            return p
        if os.path.isabs(p):
            return p
        if self.file_root:
            return os.path.join(self.file_root, p)
        return p

    def _clip(self, text: str) -> str:
        s = str(text)
        if len(s) <= self.max_chars:
            return s
        return s[: self.max_chars] + "\n... [truncated]"

    @staticmethod
    def _read_text_file(path: str) -> str:
        with open(path, "r", encoding="utf-8", errors="ignore") as fh:
            return fh.read()

    def _read_json_file(self, path: str) -> str:
        try:
            with open(path, "r", encoding="utf-8", errors="ignore") as fh:
                obj = json.load(fh)
            return json.dumps(obj, ensure_ascii=False, indent=2)
        except Exception:
            return self._read_text_file(path)

    def _read_csv_tsv(self, path: str, delimiter: str) -> str:
        rows: List[str] = []
        with open(path, "r", encoding="utf-8", errors="ignore", newline="") as fh:
            reader = csv.reader(fh, delimiter=delimiter)
            for row_idx, row in enumerate(reader):
                if row_idx >= 200:
                    rows.append("... [truncated rows]")
                    break
                rows.append("\t".join(row[:80]))
        return "\n".join(rows)

    def _read_html_file(self, path: str) -> str:
        return _strip_html(self._read_text_file(path))

    def _xml_schema_tree(self, path: str) -> str:
        tree = ET.parse(path)
        root = tree.getroot()
        lines: List[str] = [f"## Root: {root.tag}"]

        def walk(node: ET.Element, prefix: str, is_last: bool) -> None:
            children = list(node)
            seen = set()
            uniq_children = []
            for c in children:
                if c.tag in seen:
                    continue
                seen.add(c.tag)
                uniq_children.append(c)

            for i, c in enumerate(uniq_children):
                last_child = i == len(uniq_children) - 1
                connector = "└── " if last_child else "├── "
                lines.append(f"{prefix}{connector}{c.tag}")
                child_prefix = prefix + ("    " if last_child else "│   ")
                walk(c, child_prefix, last_child)

        walk(root, "", True)
        return "\n".join(lines)

    def _read_xml_file(self, path: str) -> str:
        text = self._read_text_file(path)
        if len(text) <= self.max_chars:
            return text
        try:
            return self._xml_schema_tree(path)
        except Exception:
            return text

    def _read_pdf_file(self, path: str) -> str:
        # DeepResearch also routes PDFs through a dedicated parser/IDP path.
        # Here we use local extractors when available.
        reader = None
        try:
            from pypdf import PdfReader  # type: ignore

            reader = PdfReader(path)
        except Exception:
            try:
                from PyPDF2 import PdfReader  # type: ignore

                reader = PdfReader(path)
            except Exception:
                reader = None

        if reader is None:
            return (
                "[parse_file] PDF detected but no local PDF parser available. "
                "Install `pypdf` (or `PyPDF2`) for text extraction."
            )

        pages: List[str] = []
        for idx, page in enumerate(reader.pages):
            if idx >= 80:
                pages.append("... [truncated pages]")
                break
            try:
                txt = page.extract_text() or ""
            except Exception:
                txt = ""
            if txt.strip():
                pages.append(f"[Page {idx + 1}]\n{txt.strip()}")
        if not pages:
            return "[parse_file] PDF parsed but no extractable text was found."
        return "\n\n".join(pages)

    def _read_docx_file(self, path: str) -> str:
        try:
            from docx import Document  # type: ignore

            doc = Document(path)
            chunks: List[str] = []
            for p in doc.paragraphs:
                t = (p.text or "").strip()
                if t:
                    chunks.append(t)
            for table in doc.tables:
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        chunks.append(" | ".join(cells))
            if chunks:
                return "\n".join(chunks)
        except Exception:
            pass

        # Fallback without extra deps: parse OOXML directly.
        try:
            with zipfile.ZipFile(path, "r") as zf:
                xml_data = zf.read("word/document.xml")
            root = ET.fromstring(xml_data)
            ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
            paragraphs: List[str] = []
            for p in root.findall(".//w:p", ns):
                texts = [t.text for t in p.findall(".//w:t", ns) if t.text]
                if texts:
                    paragraphs.append("".join(texts))
            return "\n".join(paragraphs)
        except Exception as exc:
            return f"[parse_file] failed to parse DOCX {path}: {exc}"

    def _read_pptx_file(self, path: str) -> str:
        try:
            from pptx import Presentation  # type: ignore

            prs = Presentation(path)
            chunks: List[str] = []
            for slide_idx, slide in enumerate(prs.slides, start=1):
                chunks.append(f"[Slide {slide_idx}]")
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False):
                        txt = (shape.text or "").strip()
                        if txt:
                            chunks.append(txt)
            if chunks:
                return "\n".join(chunks)
        except Exception:
            pass

        try:
            with zipfile.ZipFile(path, "r") as zf:
                slide_files = sorted(
                    [
                        name
                        for name in zf.namelist()
                        if name.startswith("ppt/slides/slide") and name.endswith(".xml")
                    ],
                    key=lambda s: int(re.search(r"slide(\\d+)\\.xml$", s).group(1))  # type: ignore[union-attr]
                    if re.search(r"slide(\\d+)\\.xml$", s)
                    else 10**9,
                )
                chunks: List[str] = []
                for idx, slide_file in enumerate(slide_files, start=1):
                    root = ET.fromstring(zf.read(slide_file))
                    texts = [x.text for x in root.findall(".//{*}t") if x.text]
                    if texts:
                        chunks.append(f"[Slide {idx}]")
                        chunks.extend(texts)
            if chunks:
                return "\n".join(chunks)
            return "[parse_file] PPTX parsed but no extractable text was found."
        except Exception as exc:
            return f"[parse_file] failed to parse PPTX {path}: {exc}"

    @staticmethod
    def _xlsx_shared_strings(zf: zipfile.ZipFile) -> List[str]:
        try:
            root = ET.fromstring(zf.read("xl/sharedStrings.xml"))
        except Exception:
            return []
        strings: List[str] = []
        for si in root.findall(".//{*}si"):
            text = "".join([t.text or "" for t in si.findall(".//{*}t")]).strip()
            strings.append(text)
        return strings

    @staticmethod
    def _normalize_openpyxl_color(color: Any) -> str:
        if color is None:
            return ""

        def _safe_getattr(obj: Any, name: str) -> Any:
            try:
                return getattr(obj, name)
            except Exception:
                return None

        rgb = _safe_getattr(color, "rgb")
        if isinstance(rgb, str):
            c = rgb.strip().upper()
            if len(c) == 8:  # ARGB
                c = c[2:]
            if len(c) == 6 and all(ch in "0123456789ABCDEF" for ch in c):
                return f"#{c}"

        indexed = _safe_getattr(color, "indexed")
        if isinstance(indexed, int):
            return f"indexed:{indexed}"

        theme = _safe_getattr(color, "theme")
        if theme is not None:
            tint = _safe_getattr(color, "tint")
            if tint is None:
                return f"theme:{theme}"
            try:
                return f"theme:{theme},tint:{float(tint):.2f}"
            except Exception:
                return f"theme:{theme}"
        return ""

    def _read_xlsx_color_map(self, path: str) -> str:
        """
        Extract color-layout information for style-driven spreadsheets where
        cell values may be sparse/empty (common in GAIA visual grid tasks).
        """
        try:
            from openpyxl import load_workbook  # type: ignore
        except Exception:
            return ""

        try:
            wb = load_workbook(path, data_only=True)
        except Exception:
            return ""

        sections: List[str] = []
        for ws in wb.worksheets[:4]:
            max_row = min(int(getattr(ws, "max_row", 0) or 0), 120)
            max_col = min(int(getattr(ws, "max_column", 0) or 0), 80)
            if max_row <= 0 or max_col <= 0:
                continue

            grid: List[List[str]] = []
            values: List[str] = []
            used_color = False

            for r in range(1, max_row + 1):
                row_codes: List[str] = []
                for c in range(1, max_col + 1):
                    cell = ws.cell(row=r, column=c)
                    code = ""
                    fill = getattr(cell, "fill", None)
                    if fill is not None and getattr(fill, "fill_type", None) not in (None, "none"):
                        color = self._normalize_openpyxl_color(
                            getattr(fill, "fgColor", None) or getattr(fill, "start_color", None)
                        )
                        if color:
                            code = color
                        else:
                            code = str(getattr(fill, "fill_type", "") or "").strip()

                    if code:
                        used_color = True
                    row_codes.append(code or ".")

                    val = cell.value
                    if val is not None and str(val).strip():
                        values.append(f"{cell.coordinate}={str(val).strip()}")

                while row_codes and row_codes[-1] == ".":
                    row_codes.pop()
                if row_codes:
                    grid.append(row_codes)

            if not used_color:
                continue

            colored_rows = [row for row in grid if any(x != "." for x in row)]
            if not colored_rows:
                continue

            colors = sorted({x for row in colored_rows for x in row if x != "."})
            token_map: Dict[str, str] = {color: f"C{i+1}" for i, color in enumerate(colors)}

            sections.append(f"[Sheet: {ws.title}] color map")
            sections.append("Legend: . = no fill")
            for color in colors[:20]:
                sections.append(f"{token_map[color]}={color}")
            if len(colors) > 20:
                sections.append("... [truncated legend]")

            for row_idx, row in enumerate(colored_rows[:80], start=1):
                tokens = [token_map.get(x, ".") if x != "." else "." for x in row[:80]]
                sections.append(f"R{row_idx:03d}: {' '.join(tokens)}")
            if len(colored_rows) > 80:
                sections.append("... [truncated color rows]")

            if values:
                sections.append("Cell values (non-empty):")
                for item in values[:60]:
                    sections.append(f"- {item}")
                if len(values) > 60:
                    sections.append("... [truncated values]")

        return "\n".join(sections).strip()

    def _read_xlsx_file(self, path: str) -> str:
        color_map = self._read_xlsx_color_map(path)
        try:
            import pandas as pd  # type: ignore

            xls = pd.ExcelFile(path)
            blocks: List[str] = []
            for sheet_name in xls.sheet_names[:8]:
                df = pd.read_excel(path, sheet_name=sheet_name, dtype=str).fillna("")
                if df.empty:
                    blocks.append(f"[Sheet: {sheet_name}] (empty)")
                    continue
                non_empty_mask = df.apply(lambda r: any(str(x).strip() for x in r), axis=1)
                df = df[non_empty_mask]
                if df.empty:
                    blocks.append(f"[Sheet: {sheet_name}] (only empty rows)")
                    continue
                blocks.append(f"[Sheet: {sheet_name}]")
                blocks.append(df.head(120).to_csv(index=False, sep="\t"))
            if color_map:
                blocks.append(color_map)
            if blocks:
                return "\n".join(blocks)
        except Exception:
            pass

        # Fallback without pandas/openpyxl.
        try:
            with zipfile.ZipFile(path, "r") as zf:
                shared = self._xlsx_shared_strings(zf)
                sheet_files = sorted(
                    [
                        name
                        for name in zf.namelist()
                        if name.startswith("xl/worksheets/sheet") and name.endswith(".xml")
                    ],
                    key=lambda s: int(re.search(r"sheet(\\d+)\\.xml$", s).group(1))  # type: ignore[union-attr]
                    if re.search(r"sheet(\\d+)\\.xml$", s)
                    else 10**9,
                )
                blocks: List[str] = []
                for sheet_idx, sheet_file in enumerate(sheet_files[:8], start=1):
                    root = ET.fromstring(zf.read(sheet_file))
                    blocks.append(f"[Sheet {sheet_idx}]")
                    row_count = 0
                    for row in root.findall(".//{*}sheetData/{*}row"):
                        row_count += 1
                        if row_count > 200:
                            blocks.append("... [truncated rows]")
                            break
                        vals: List[str] = []
                        for cell in row.findall("{*}c"):
                            ctype = cell.attrib.get("t", "")
                            value = ""
                            v = cell.find("{*}v")
                            if ctype == "inlineStr":
                                t = cell.find("{*}is/{*}t")
                                value = (t.text or "") if t is not None else ""
                            elif ctype == "s" and v is not None and v.text is not None:
                                try:
                                    idx = int(v.text)
                                    value = shared[idx] if 0 <= idx < len(shared) else v.text
                                except Exception:
                                    value = v.text or ""
                            else:
                                value = (v.text or "") if v is not None else ""
                            vals.append(value)
                        if vals:
                            blocks.append("\t".join(vals))
                if color_map:
                    blocks.append(color_map)
                if blocks:
                    return "\n".join(blocks)
        except Exception as exc:
            return f"[parse_file] failed to parse XLSX {path}: {exc}"
        if color_map:
            return color_map
        return "[parse_file] XLSX parsed but no extractable rows were found."

    def _read_xls_file(self, path: str) -> str:
        try:
            import pandas as pd  # type: ignore

            df = pd.read_excel(path, dtype=str).fillna("")
            return df.head(120).to_csv(index=False, sep="\t")
        except Exception:
            return (
                "[parse_file] XLS detected. Install `pandas` + `xlrd` for robust "
                "legacy XLS parsing, or convert to XLSX/CSV."
            )

    def _read_media_file(self, path: str) -> str:
        ffprobe_bin = shutil.which("ffprobe")
        if ffprobe_bin:
            try:
                proc = subprocess.run(
                    [
                        ffprobe_bin,
                        "-v",
                        "error",
                        "-show_entries",
                        "format=filename,format_name,duration,size,bit_rate",
                        "-show_streams",
                        "-of",
                        "json",
                        path,
                    ],
                    capture_output=True,
                    text=True,
                    timeout=25,
                )
                if proc.returncode == 0 and proc.stdout.strip():
                    payload = json.loads(proc.stdout)
                    return json.dumps(payload, ensure_ascii=False, indent=2)
            except Exception:
                pass
        return (
            "[parse_file] media file detected. Metadata extraction requires `ffprobe` "
            "(from FFmpeg). Full transcription is not supported in this local tool."
        )

    def _read_image_file(self, path: str) -> str:
        lines: List[str] = []
        try:
            from PIL import Image  # type: ignore

            image = Image.open(path)
            lines.append(
                f"[image] format={image.format or 'unknown'}, size={image.size[0]}x{image.size[1]}, mode={image.mode}"
            )
            try:
                import pytesseract  # type: ignore

                ocr_text = (pytesseract.image_to_string(image) or "").strip()
                if ocr_text:
                    lines.append("[ocr]")
                    lines.append(ocr_text)
                else:
                    lines.append("[ocr] no text detected")
            except Exception:
                lines.append(
                    "[parse_file] OCR not available. Install `pytesseract` and system `tesseract` binary for text extraction."
                )
            return "\n".join(lines)
        except Exception:
            return "[parse_file] image parsing requires Pillow (`pip install pillow`)."

    def _read_zip_file(self, path: str, *, depth: int) -> str:
        if depth >= 2:
            return "[parse_file] ZIP nesting limit reached."
        sections: List[str] = []
        with tempfile.TemporaryDirectory(prefix="parse_file_zip_") as tmpdir:
            with zipfile.ZipFile(path, "r") as zf:
                names = [n for n in zf.namelist() if not n.endswith("/")]
                if not names:
                    return "[parse_file] ZIP is empty."
                for name in names[:25]:
                    try:
                        extracted = zf.extract(name, path=tmpdir)
                    except Exception as exc:
                        sections.append(f"[zip] failed extracting {name}: {exc}")
                        continue
                    body = self._read_one_internal(extracted, depth=depth + 1)
                    sections.append(f"## {name}\n{body}")
                if len(names) > 25:
                    sections.append("... [truncated zip entries]")
        return "\n\n".join(sections)

    def _read_one_internal(self, resolved: str, *, depth: int = 0) -> str:
        if not os.path.exists(resolved):
            return f"[parse_file] missing: {resolved}"

        ext = Path(resolved).suffix.lower()
        if ext in {".txt", ".md", ".py", ".jsonl", ".yaml", ".yml", ".ini", ".cfg", ".log"}:
            return self._read_text_file(resolved)
        if ext == ".json":
            return self._read_json_file(resolved)
        if ext == ".csv":
            return self._read_csv_tsv(resolved, ",")
        if ext == ".tsv":
            return self._read_csv_tsv(resolved, "\t")
        if ext in {".html", ".htm"}:
            return self._read_html_file(resolved)
        if ext == ".xml":
            return self._read_xml_file(resolved)
        if ext == ".pdf":
            return self._read_pdf_file(resolved)
        if ext in {".docx", ".doc"}:
            return self._read_docx_file(resolved)
        if ext == ".pptx":
            return self._read_pptx_file(resolved)
        if ext == ".xlsx":
            return self._read_xlsx_file(resolved)
        if ext == ".xls":
            return self._read_xls_file(resolved)
        if ext in {".png", ".jpg", ".jpeg", ".bmp", ".gif", ".webp", ".tiff", ".tif"}:
            return self._read_image_file(resolved)
        if ext in {".mp3", ".wav", ".aac", ".ogg", ".flac", ".mp4", ".mov", ".avi", ".mkv", ".webm"}:
            return self._read_media_file(resolved)
        if ext == ".zip":
            return self._read_zip_file(resolved, depth=depth)

        mime, _ = mimetypes.guess_type(resolved)
        if mime and mime.startswith("text/"):
            return self._read_text_file(resolved)
        return f"[parse_file] unsupported extension: {ext or 'unknown'}"

    def _read_one(self, path: str) -> str:
        resolved = self._resolve(path)
        try:
            body = self._read_one_internal(resolved, depth=0)
            return f"# File: {resolved}\n{self._clip(body)}"
        except Exception as exc:
            return f"[parse_file] failed to read {resolved}: {exc}"

    def call(self, arguments: Dict[str, Any]) -> str:
        files = arguments.get("files")
        if isinstance(files, str):
            files = [files]
        if not isinstance(files, list):
            return "[parse_file] invalid arguments: files must be list of paths"

        outputs = [self._read_one(str(p)) for p in files[:10]]
        return "\n\n=======\n\n".join(outputs)


class ToolRegistry:
    def __init__(self, tools: List[BaseTool]) -> None:
        self._tools = {tool.spec.name: tool for tool in tools}

    def schemas(self) -> List[Dict[str, Any]]:
        out = []
        for tool in self._tools.values():
            out.append(
                {
                    "type": "function",
                    "function": {
                        "name": tool.spec.name,
                        "description": tool.spec.description,
                        "parameters": tool.spec.parameters,
                    },
                }
            )
        return out

    def call(self, name: str, arguments: Dict[str, Any]) -> str:
        tool = self._tools.get(name)
        if tool is None:
            return f"Error: tool '{name}' not found"
        return tool.call(arguments)



def build_default_tool_registry(
    file_root: str = "",
    *,
    enable_python_interpreter: bool = True,
) -> ToolRegistry:
    tools: List[BaseTool] = [
        SearchTool(),
        VisitTool(),
        ParseFileTool(file_root=file_root),
    ]
    if enable_python_interpreter:
        tools.insert(2, PythonInterpreterTool())
    return ToolRegistry(tools)
